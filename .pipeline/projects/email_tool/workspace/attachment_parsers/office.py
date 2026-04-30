"""Office document attachment parser (DOCX, XLSX, PPTX)."""

import os
from typing import Optional

from email_tool.attachment_parsers.base import ParsedAttachment, AttachmentMetadata


class OfficeAttachmentParser:
    """Parser for Office document attachments."""
    
    def __init__(self, base_path: str):
        """Initialize the Office parser."""
        self.base_path = base_path
        self.parser_name = "OfficeAttachmentParser"
    
    def parse(
        self,
        file_path: str,
        attachment_id: str,
        email_id: str,
        original_filename: str,
        content_type: str,
        size_bytes: int,
        attachment_type: str
    ) -> ParsedAttachment:
        """Parse an Office document and extract text content."""
        try:
            if not os.path.exists(file_path):
                return ParsedAttachment(
                    success=False,
                    error_message=f"File not found: {file_path}",
                    attachment_id=attachment_id,
                    email_id=email_id,
                    original_filename=original_filename,
                    content_type=content_type,
                    size_bytes=size_bytes,
                    attachment_type=attachment_type
                )
            
            text_content = ""
            
            if original_filename.lower().endswith('.docx'):
                # Extract text from DOCX
                try:
                    from docx import Document
                    doc = Document(file_path)
                    text_parts = [para.text for para in doc.paragraphs]
                    text_content = "\n".join(text_parts)
                except Exception as e:
                    text_content = f"[DOCX extraction error: {str(e)}]"
            
            elif original_filename.lower().endswith('.xlsx'):
                # Extract text from XLSX
                try:
                    import pandas as pd
                    excel_file = pd.ExcelFile(file_path)
                    text_parts = []
                    for sheet_name in excel_file.sheet_names:
                        df = pd.read_excel(excel_file, sheet_name=sheet_name)
                        text_parts.append(f"Sheet: {sheet_name}\n{df.to_string()}")
                    text_content = "\n\n".join(text_parts)
                except Exception as e:
                    text_content = f"[XLSX extraction error: {str(e)}]"
            
            elif original_filename.lower().endswith('.pptx'):
                # Extract text from PPTX
                try:
                    from pptx import Presentation
                    prs = Presentation(file_path)
                    text_parts = []
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                text_parts.append(shape.text)
                    text_content = "\n".join(text_parts)
                except Exception as e:
                    text_content = f"[PPTX extraction error: {str(e)}]"
            
            metadata = AttachmentMetadata(
                file_size=size_bytes,
                file_type=content_type,
                parser_name=self.parser_name
            )
            
            return ParsedAttachment(
                success=True,
                text_content=text_content,
                metadata=metadata,
                attachment_id=attachment_id,
                email_id=email_id,
                original_filename=original_filename,
                content_type=content_type,
                size_bytes=size_bytes,
                attachment_type=attachment_type
            )
            
        except Exception as e:
            return ParsedAttachment(
                success=False,
                error_message=str(e),
                attachment_id=attachment_id,
                email_id=email_id,
                original_filename=original_filename,
                content_type=content_type,
                size_bytes=size_bytes,
                attachment_type=attachment_type
            )
